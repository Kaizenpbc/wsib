from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import Dict, Any, Optional
import os
from sqlalchemy.orm import Session

class PptGenerator:
    """
    Service to generate PowerPoint presentations from curriculum
    Now reads actual curriculum from database!
    """
    
    def __init__(self, db: Optional[Session] = None):
        self.db = db
    
    async def generate(self, curriculum_id: str, customizations: Dict[str, Any] = None) -> str:
        """
        Generate a PowerPoint presentation from curriculum data
        
        Args:
            curriculum_id: ID of curriculum to generate from
            customizations: Optional dict with theme, colors, preferences
        """
        from app.db_models import Curriculum
        
        # Get curriculum from database
        if self.db:
            curriculum = self.db.query(Curriculum).filter(Curriculum.id == curriculum_id).first()
        else:
            curriculum = None
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Get customizations
        theme = customizations.get('theme', 'professional') if customizations else 'professional'
        color_scheme = customizations.get('color_scheme', 'blue') if customizations else 'blue'
        
        # Title
        title = curriculum.title if curriculum else "Training Program"
        subtitle = curriculum.description if curriculum else "Professional Curriculum"
        
        # Add title slide
        self._add_title_slide(prs, title, subtitle)
        
        # If we have curriculum data, use it
        if curriculum and hasattr(curriculum, 'total_duration_minutes'):
            # Add overview slide
            self._add_content_slide(
                prs,
                "Course Overview",
                [
                    f"Duration: {curriculum.total_duration_minutes // 60} hours",
                    f"Status: {curriculum.status.title()}",
                    "Comprehensive training program",
                    "Hands-on practice included"
                ]
            )
        
        # Add generic content slides (will be enhanced with actual modules)
        self._add_content_slide(
            prs,
            "Learning Objectives",
            [
                "Understand key concepts and principles",
                "Demonstrate practical skills",
                "Apply knowledge in real scenarios",
                "Pass competency assessments"
            ]
        )
        
        self._add_content_slide(
            prs,
            "Course Content",
            [
                "Comprehensive curriculum coverage",
                "2. Call for help",
                "3. Position the patient",
                "4. Begin chest compressions",
                "5. Give rescue breaths"
            ]
        )
        
        # Save file
        output_dir = "outputs/presentations"
        os.makedirs(output_dir, exist_ok=True)
        file_path = f"{output_dir}/curriculum_{curriculum_id}.pptx"
        prs.save(file_path)
        
        return file_path
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """Add a title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
    
    def _add_content_slide(self, prs: Presentation, title: str, bullet_points: list):
        """Add a content slide with bullet points"""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)

